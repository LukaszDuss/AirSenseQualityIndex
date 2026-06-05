# -*- coding: utf-8 -*-
"""Buduje prezentację AirSense Weather AI (16:9) za pomocą python-pptx."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "Prezentacja_AirSense_Weather_AI.pptx")
def A(n): return os.path.join(ASSETS, n)

# --- Paleta (spójna z aplikacją) ---
NAVY  = RGBColor(0x0E, 0x22, 0x38)
NAVY2 = RGBColor(0x14, 0x30, 0x52)
BLUE  = RGBColor(0x3B, 0x82, 0xF6)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
INK   = RGBColor(0x12, 0x26, 0x3D)
MUTED = RGBColor(0x5C, 0x6B, 0x80)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
CLINE = RGBColor(0xD8, 0xE3, 0xF0)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ICEBL = RGBColor(0xCA, 0xDC, 0xFC)

HEAD = "Trebuchet MS"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ----------------- helpers -----------------
def slide(bg=LIGHT):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def _box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf

def _run(p, text, size, color, bold=False, font=BODY, italic=False):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    return r

def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    return p

def line_text(s, l, t, w, h, text, size, color, bold=False, font=BODY,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb, tf = _box(s, l, t, w, h); tf.vertical_anchor = anchor
    _run(para(tf, first=True, align=align), text, size, color, bold, font, italic)
    return tb

def rrect(s, l, t, w, h, fill=CARD, line=CLINE, radius=0.06):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    try: sh.adjustments[0] = radius
    except Exception: pass
    return sh

def circle(s, l, t, d, fill, text=None, tcolor=NAVY, tsize=18):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    sh.shadow.inherit = False
    if text is not None:
        tf = sh.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(para(tf, first=True, align=PP_ALIGN.CENTER), text, tsize, tcolor, True, HEAD)
    return sh

def vdash(s, x, y, h, color=AMBER, width=2.0):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x), Inches(y + h))
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    d = ln.find(qn('a:prstDash'))
    if d is None:
        d = ln.makeelement(qn('a:prstDash'), {}); ln.append(d)
    d.set('val', 'dash')
    cn.shadow.inherit = False
    return cn

def pic(s, name, l, t, w=None, h=None):
    kw = {}
    if w is not None: kw['width'] = Inches(w)
    if h is not None: kw['height'] = Inches(h)
    return s.shapes.add_picture(A(name), Inches(l), Inches(t), **kw)

def title_block(s, num, title, num_color=AMBER, title_color=INK):
    circle(s, 0.6, 0.52, 0.66, num_color, str(num), NAVY, 24)
    line_text(s, 1.46, 0.5, 11.3, 0.92, title, 30, title_color, True, HEAD,
              anchor=MSO_ANCHOR.MIDDLE)

def footer(s, idx):
    line_text(s, 0.6, 7.04, 9.0, 0.34, "AirSense Weather AI   ·   Grupa 5", 10, MUTED)
    line_text(s, 11.4, 7.04, 1.33, 0.34, str(idx), 10, MUTED, align=PP_ALIGN.RIGHT)

def bullets(s, l, t, w, h, items, size=15, gap=10, color=INK, marker_color=BLUE):
    tb, tf = _box(s, l, t, w, h)
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0), space_after=gap)
        _run(p, "•  ", size, marker_color, True)
        if isinstance(it, tuple):  # (lead_bold, rest)
            _run(p, it[0], size, color, True)
            _run(p, it[1], size, color, False)
        else:
            _run(p, it, size, color, False)
    return tb


# ===================== SLIDE 1 — TITLE =====================
s = slide(NAVY)
# motyw: pionowa przerywana linia "Dziś" (sygnatura projektu)
vdash(s, 10.7, 1.2, 5.0, AMBER, 2.2)
line_text(s, 10.2, 0.85, 2.2, 0.4, "Dziś", 14, AMBER, True, HEAD, align=PP_ALIGN.CENTER)
line_text(s, 0.9, 2.05, 10.5, 1.2, "AirSense Weather AI", 54, WHITE, True, HEAD)
line_text(s, 0.95, 3.25, 10.0, 0.7, "Inteligentny system predykcji jakości powietrza i pogody",
          23, ICEBL, False, BODY)
line_text(s, 0.95, 4.05, 10.5, 0.5, "Sekwencyjny model LSTM  ·  Streamlit  ·  dane GIOŚ + Open-Meteo",
          15, AMBER, False, BODY)
line_text(s, 0.95, 5.7, 11.0, 0.5, "Grupa 5", 18, WHITE, True, HEAD)
line_text(s, 0.95, 6.15, 11.4, 0.5,
          "Jakub Moszyński  ·  Damian Tomczyk  ·  Nikol Olszewska  ·  Łukasz Duss",
          14, ICEBL, False, BODY)
notes(s, "Powitanie. Przedstawienie tematu: prototyp systemu, który łączy realne dane pogodowe (Open-Meteo) "
         "i jakości powietrza (GIOŚ) i prognozuje je na kolejne dni siecią LSTM. Przedstawienie zespołu. (~1 min)")

# ===================== SLIDE 2 — PLAN =====================
s = slide(LIGHT)
title_block(s, "•", "Plan prezentacji")
agenda = [
    ("1", "Przegląd projektu i problem"),
    ("2", "Dane i wyzwania preprocessingu"),
    ("3", "Model AI i jego uzasadnienie"),
    ("4", "Wyniki, metryki i wizualizacje"),
    ("5", "Demo aplikacji"),
    ("6", "Wnioski i kierunki rozwoju"),
]
cw, ch, gx, gy = 5.85, 1.28, 0.42, 0.34
x0, y0 = 0.6, 1.95
for i, (n, label) in enumerate(agenda):
    col, row = i % 2, i // 2
    l = x0 + col * (cw + gx); t = y0 + row * (ch + gy)
    rrect(s, l, t, cw, ch)
    circle(s, l + 0.3, t + (ch - 0.62) / 2, 0.62, AMBER, n, NAVY, 22)
    line_text(s, l + 1.12, t, cw - 1.3, ch, label, 17, INK, True, HEAD, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)
notes(s, "Krótki przegląd struktury wystąpienia — 6 punktów wymaganych na ocenę. (~30 s)")

# ===================== SLIDE 3 — ZESPÓŁ =====================
s = slide(LIGHT)
title_block(s, "•", "Zespół i podział pracy")
team = [
    ("Jakub Moszyński", "Data Engineer", "Integracja API GIOŚ i Open-Meteo, pobieranie i czyszczenie danych, przygotowanie sekwencji czasowych.", BLUE),
    ("Damian Tomczyk", "Data Analyst", "Analiza eksploracyjna (EDA), dashboard Streamlit, wizualizacje trendów w Plotly, linia „Dziś”.", GREEN),
    ("Nikol Olszewska", "Deep Learning Engineer", "Architektura LSTM Seq2Seq, Dropout, skalowanie danych, dynamiczna liczba cech.", RGBColor(0x7C,0x3A,0xED)),
    ("Łukasz Duss", "ML Analyst / Project Manager", "Integracja aplikacji, podział 80/10/10, ewaluacja MAE/MSE/RMSE/R², rejestr modeli, Session State.", AMBER),
]
cw, ch, gx, gy = 5.85, 2.18, 0.42, 0.3
x0, y0 = 0.6, 1.95
for i, (name, role, contrib, acc) in enumerate(team):
    col, row = i % 2, i // 2
    l = x0 + col * (cw + gx); t = y0 + row * (ch + gy)
    rrect(s, l, t, cw, ch)
    circle(s, l + 0.32, t + 0.34, 0.34, acc)
    line_text(s, l + 0.86, t + 0.26, cw - 1.1, 0.5, name, 17, INK, True, HEAD)
    line_text(s, l + 0.86, t + 0.74, cw - 1.1, 0.4, role, 12.5, acc, True, BODY)
    tb, tf = _box(s, l + 0.32, t + 1.22, cw - 0.64, ch - 1.35)
    _run(para(tf, first=True, space_after=0), contrib, 12.5, MUTED)
footer(s, 3)
notes(s, "Kto za co odpowiadał. Każda osoba przedstawia w demie swoją część. Szczegóły w raporcie indywidualnym. (~1 min)")

# ===================== SLIDE 4 — PROBLEM =====================
s = slide(LIGHT)
title_block(s, 1, "Problem: jakość powietrza i prognoza")
bullets(s, 0.65, 1.95, 7.1, 4.6, [
    "Polska zmaga się z wysokim poziomem smogu (PM10, PM2.5, NO2), szczególnie w sezonie grzewczym.",
    "Decyzje (aktywność na zewnątrz, ostrzeżenia) wymagają prognozy na kolejne dni — nie tylko bieżącego odczytu.",
    "Dane są rozproszone: pogoda w jednym źródle, zanieczyszczenia w drugim.",
    "Potrzebny jest system, który łączy te źródła i automatycznie prognozuje kolejne dni.",
], size=16, gap=14)
# karta celu
rrect(s, 8.05, 1.95, 4.68, 2.5, fill=NAVY, line=None)
line_text(s, 8.4, 2.2, 4.0, 0.5, "Nasz cel", 18, AMBER, True, HEAD)
tb, tf = _box(s, 8.4, 2.78, 4.0, 1.55)
_run(para(tf, first=True), "Prototyp łączący realne dane pogodowe i środowiskowe i prognozujący je krótkoterminowo (godzinowo) za pomocą sieci LSTM.",
     14.5, WHITE)
# 3 staty
for i, (num, lab) in enumerate([("5", "stacji pomiarowych"), ("+1–24", "godziny prognozy"), ("9", "śledzonych cech")]):
    l = 8.05 + i * 1.62
    line_text(s, l, 4.75, 1.55, 0.7, num, 34, BLUE, True, HEAD)
    line_text(s, l, 5.5, 1.55, 0.7, lab, 11.5, MUTED)
footer(s, 4)
notes(s, "Definicja problemu: smog to realny problem zdrowotny. Brak jednego narzędzia łączącego pogodę i jakość "
         "powietrza z prognozą. Nasz cel — prototyp end-to-end. (~2 min)")

# ===================== SLIDE 5 — PIPELINE =====================
s = slide(LIGHT)
title_block(s, 1, "Pełny pipeline Machine Learning")
pic(s, "pipeline.png", 0.6, 2.55, w=12.13)
line_text(s, 0.65, 5.55, 12.0, 1.0,
          "Od pobrania danych z publicznych API, przez czyszczenie i sekwencje czasowe, po model LSTM, "
          "ewaluację i predykcję — wszystko w jednej aplikacji Streamlit. Projekt obejmuje pełny cykl ML "
          "(prototyp klasy MLOps).", 15, INK)
footer(s, 5)
notes(s, "Pokazujemy, że to nie tylko model, ale cały pipeline: dane → preprocessing → sekwencje → LSTM → ewaluacja → predykcja. "
         "Każdy etap omówimy na kolejnych slajdach. (~1 min)")

# ===================== SLIDE 6 — ŹRÓDŁA DANYCH =====================
s = slide(LIGHT)
title_block(s, 2, "Skąd pochodzą dane")
# karta Open-Meteo
rrect(s, 0.6, 1.95, 5.95, 3.5)
circle(s, 0.9, 2.25, 0.34, BLUE)
line_text(s, 1.4, 2.18, 4.9, 0.5, "Open-Meteo — pogoda", 18, INK, True, HEAD)
bullets(s, 0.95, 2.95, 5.3, 2.4, [
    "Temperatura, wilgotność, prędkość wiatru",
    "Dane godzinowe (~720 punktów na 30 dni)",
    "forecast_days=1 — bez godzin „w przyszłość”",
], size=13.5, gap=9)
# karta GIOŚ
rrect(s, 6.78, 1.95, 5.95, 3.5)
circle(s, 7.08, 2.25, 0.34, GREEN)
line_text(s, 7.58, 2.18, 4.9, 0.5, "GIOŚ API v1 — zanieczyszczenia", 18, INK, True, HEAD)
bullets(s, 7.13, 2.95, 5.3, 2.4, [
    "PM10, PM2.5, NO2, NOx, O3, NO — zależnie od stacji",
    "Nagłówki HTTP imitujące przeglądarkę (omijają blokady)",
    "Uwaga: API daje tylko ~3 dni historii godzinowej",
], size=13.5, gap=9, marker_color=GREEN)
line_text(s, 0.6, 5.75, 12.1, 0.6,
          "5 stacji pomiarowych:   Zabrze  ·  Warszawa  ·  Kraków  ·  Gdańsk  ·  Wrocław", 15, INK, True, BODY)
footer(s, 6)
notes(s, "Dwa źródła. Open-Meteo: pogoda, dane godzinowe (~720 pkt), 30 dni wstecz, pobieramy tylko do teraz. "
         "GIOŚ: zanieczyszczenia, v1 API, różne sensory na stacjach, wymaga nagłówków HTTP. (~1.5 min)")

# ===================== SLIDE 7 — WYZWANIA PREPROCESSINGU =====================
s = slide(LIGHT)
title_block(s, 2, "Wyzwania w przygotowaniu danych")
rows = [
    ("Braki w pomiarach", "Interpolacja liniowa w obie strony (limit_direction='both').", BLUE),
    ("Błędne piki czujników", "Przycinanie do 1. i 99. percentyla (winsoryzacja).", GREEN),
    ("Różna liczba sensorów", "Dynamiczna liczba cech wyznaczana per stacja (CURRENT_FEATURES).", RGBColor(0x7C,0x3A,0xED)),
    ("Blokady API GIOŚ", "Nagłówki HTTP (User-Agent, Referer) + obsługa błędów.", RGBColor(0x0E,0xA5,0xB5)),
    ("Dni „w przyszłość” z czujek", "forecast_days=1 + obcięcie do dziś; pionowa linia „Dziś” na wykresie.", AMBER),
]
y = 1.85
for head, desc, acc in rows:
    circle(s, 0.72, y + 0.08, 0.5, acc)
    line_text(s, 1.45, y, 11.2, 0.45, head, 16.5, INK, True, HEAD)
    line_text(s, 1.45, y + 0.45, 11.2, 0.5, desc, 13.5, MUTED)
    y += 1.0
footer(s, 7)
notes(s, "Realne dane bywają brudne. Pokazujemy konkretne problemy i rozwiązania. Ostatni punkt to nasza ostatnia poprawka: "
         "nie pobieramy dni w przyszłość i dodaliśmy linię „Dziś”. (~2 min)")

# ===================== SLIDE 8 — ARCHITEKTURA =====================
s = slide(LIGHT)
title_block(s, 3, "Model: LSTM Sequence-to-Sequence")
pic(s, "architecture.png", 0.6, 1.85, h=4.95)
bullets(s, 7.4, 2.2, 5.3, 4.4, [
    ("LSTM ", "dobrze modeluje zależności w szeregach czasowych (pamięć krótkiej historii)."),
    ("Seq2Seq: ", "z okna 48 h przewiduje kolejne godziny naprzód jednocześnie."),
    ("Wielowymiarowość: ", "wszystkie cechy (pogoda + zanieczyszczenia) na wejściu i wyjściu."),
    ("Dropout 0.2: ", "ogranicza przeuczenie przy stosunkowo małym zbiorze."),
    ("Dynamiczne wejście: ", "kształt dopasowany do liczby sensorów danej stacji."),
], size=15, gap=14)
footer(s, 8)
notes(s, "Dlaczego LSTM: szeregi czasowe, pamięć, sekwencja na sekwencję. Diagram po lewej: dwie warstwy LSTM(64) z "
         "Dropout, Dense + Reshape na wyjściu (horyzont × cechy). (~2 min)")

# ===================== SLIDE 9 — TRENING =====================
s = slide(LIGHT)
title_block(s, 3, "Trening i konfiguracja modelu")
cfg = [
    ("Podział 80 / 10 / 10", "trening / walidacja / test — chronologicznie, bez tasowania."),
    ("Okno 48 h → horyzont +1–24 h", "model godzinowy; ~700 punktów na 30 dni."),
    ("2× LSTM(64) + Dropout", "Dense + Reshape na wyjściu sekwencyjnym."),
    ("Adam · MSE · early stopping", "MinMaxScaler; parametry skalera zapisane do predykcji."),
]
cw, ch, gx, gy = 5.85, 1.62, 0.42, 0.34
x0, y0 = 0.6, 2.0
for i, (head, desc) in enumerate(cfg):
    col, row = i % 2, i // 2
    l = x0 + col * (cw + gx); t = y0 + row * (ch + gy)
    rrect(s, l, t, cw, ch)
    line_text(s, l + 0.32, t + 0.22, cw - 0.64, 0.5, head, 16, INK, True, HEAD)
    tb, tf = _box(s, l + 0.32, t + 0.74, cw - 0.64, ch - 0.85)
    _run(para(tf, first=True), desc, 13, MUTED)
footer(s, 9)
notes(s, "Konfiguracja treningu. Split chronologiczny (ważne dla szeregów czasowych). Mały model + Dropout, bo zbiór jest "
         "krótki. Skaler zapisujemy, żeby odtworzyć skalę przy predykcji. (~1.5 min)")

# ===================== SLIDE 10 — METRYKI =====================
s = slide(LIGHT)
title_block(s, 4, "Wyniki — realne metryki (dane godzinowe)")
# lewe: uczciwe stat callouts
for i, (num, lab, c) in enumerate([("13%", "błąd wzgl. wilgotności (najlepszy)", GREEN),
                                   ("~3°C", "średni błąd temperatury (MAE)", BLUE),
                                   ("≈ 0", "R² — poziom baseline'u", AMBER)]):
    t = 1.9 + i * 1.12
    line_text(s, 0.7, t, 2.1, 0.8, num, 34, c, True, HEAD)
    line_text(s, 2.75, t + 0.12, 3.4, 0.8, lab, 12.5, MUTED, True)
rrect(s, 0.65, 5.4, 5.65, 1.5, fill=LIGHT, line=CLINE)
tb, tf = _box(s, 0.9, 5.56, 5.15, 1.25)
_run(para(tf, first=True, space_after=4), "Metryki realne, w jednostkach rzeczywistych. ", 12, INK, True)
_run(tf.paragraphs[0], "Model łapie krótkoterminowy sygnał (wilgotność, pogoda).", 12, INK)
_run(para(tf), "R² bliskie 0 i niestabilne — mały zbiór testowy i tylko ~3 dni danych GIOŚ.", 12, MUTED)
pic(s, "metrics_per_feature.png", 6.5, 2.0, w=6.3)
footer(s, 10)
notes(s, "Uczciwie: to realne wyniki z danych godzinowych (+1 h). MAE jest stabilne i interpretowalne (np. temp ~3°C, "
         "wilgotność ~13% błędu względnego). R² jest bliskie 0 / niestabilne (od ~-1 do ~+0.5 między uruchomieniami) — "
         "bo zbiór testowy to ~3 dni, a dane GIOŚ o zanieczyszczeniach mają tylko ~3 dni historii. To prototyp, główny "
         "kierunek rozwoju = więcej danych. (~2 min)")

# ===================== SLIDE 11 — WIZUALIZACJE =====================
s = slide(LIGHT)
title_block(s, 4, "Wizualizacje w aplikacji")
pic(s, "eda_trend.png", 0.6, 1.95, w=6.05)
pic(s, "prediction.png", 6.85, 1.95, w=6.05)
line_text(s, 0.6, 4.95, 6.05, 1.4,
          "EDA: dane godzinowe (cykl dobowy) + 24-godzinna średnia krocząca + pionowa linia „Dziś”.",
          13, INK)
line_text(s, 6.85, 4.95, 6.05, 1.4,
          "Predykcja vs rzeczywistość (zbiór testowy, +1 h): model śledzi ogólny kształt, ale z wyraźnymi błędami.",
          13, INK)
footer(s, 11)
notes(s, "Dwa kluczowe wykresy z aplikacji. Po lewej zakładka 1 (trend + „Dziś”), po prawej zakładka 3 (prognoza połączona z "
         "ostatnim pomiarem, prawdziwe daty na osi). (~1.5 min)")

# ===================== SLIDE 12 — DEMO =====================
s = slide(LIGHT)
title_block(s, 5, "Demo aplikacji (Streamlit)")
tabs = [
    ("Zakładka 1", "EDA i Trendy", "Wybór stacji i parametru; wykres odczytów + średnia krocząca 24h + linia „Dziś”.", BLUE),
    ("Zakładka 2", "Model LSTM", "Konfiguracja i trening modelu; podział 80/10/10; pasek postępu epok na żywo.", RGBColor(0x7C,0x3A,0xED)),
    ("Zakładka 3", "Predykcja i Ewaluacja", "Metryki, edytowalna macierz wejściowa, interaktywny wykres prognozy.", GREEN),
]
cw, gx = 3.93, 0.27
x0, t = 0.6, 2.0
for i, (tab, name, desc, acc) in enumerate(tabs):
    l = x0 + i * (cw + gx)
    rrect(s, l, t, cw, 3.5)
    circle(s, l + 0.3, t + 0.32, 0.5, acc, str(i + 1), WHITE, 18)
    line_text(s, l + 0.95, t + 0.28, cw - 1.1, 0.4, tab, 12.5, acc, True, BODY)
    line_text(s, l + 0.95, t + 0.66, cw - 1.1, 0.5, name, 16, INK, True, HEAD)
    tb, tf = _box(s, l + 0.32, t + 1.35, cw - 0.64, 1.9)
    _run(para(tf, first=True), desc, 13, MUTED)
rrect(s, 0.6, 5.85, 12.13, 0.95, fill=NAVY, line=None)
tb, tf = _box(s, 0.95, 6.02, 11.5, 0.7); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_run(para(tf, first=True), "Demo na żywo:  ", 15, AMBER, True)
_run(tf.paragraphs[0], "streamlit run app.py", 15, WHITE, True, "Consolas")
_run(tf.paragraphs[0], "    (plan awaryjny: zrzuty ekranu / nagranie)", 13, ICEBL)
footer(s, 12)
notes(s, "Pokaz na żywo: uruchamiamy aplikację, wybieramy stację, trenujemy szybki model, generujemy prognozę. "
         "Trzy zakładki = trzy etapy. Gdyby brakło sieci — zrzuty ekranu jako plan B. (~3–4 min)")

# ===================== SLIDE 13 — WNIOSKI =====================
s = slide(LIGHT)
title_block(s, 6, "Wnioski")
bullets(s, 0.7, 2.0, 12.0, 4.6, [
    ("Działający prototyp klasy MLOps: ", "pełny pipeline od publicznych API po predykcję w aplikacji webowej."),
    ("Realne dane godzinowe: ", "GIOŚ (zanieczyszczenia) + Open-Meteo (pogoda), czyszczone i łączone automatycznie."),
    ("Krótkoterminowo model łapie sygnał pogodowy ", "(wilgotność, temperatura) — ale to wciąż prototyp, nie model produkcyjny."),
    ("Główne ograniczenie: ", "tylko ~3 dni danych GIOŚ i mały zbiór testowy → R² bliskie 0 i niestabilne."),
    ("Uczciwa ewaluacja: ", "metryki w jednostkach rzeczywistych, w rozbiciu na cechy, na wydzielonym zbiorze testowym."),
], size=16, gap=15)
footer(s, 13)
notes(s, "Podsumowanie: osiągnęliśmy cel — działający, kompletny prototyp na realnych danych. Świadomi ograniczeń "
         "(dane, dokładność pollutantów). (~1.5 min)")

# ===================== SLIDE 14 — ROZWÓJ =====================
s = slide(LIGHT)
title_block(s, 6, "Kierunki dalszego rozwoju")
fut = [
    "Dłuższa historia danych niż 30 dni",
    "Więcej miast + automatyczne wykrywanie stacji",
    "Porównanie z GRU, Prophet, XGBoost, Transformerami",
    "Automatyczny, cykliczny re-trening modeli",
    "Zapis historii predykcji w bazie danych",
    "Rozbudowany rejestr modeli z porównaniem metryk",
]
cw, ch, gx, gy = 5.85, 1.28, 0.42, 0.34
x0, y0 = 0.6, 1.95
for i, item in enumerate(fut):
    col, row = i % 2, i // 2
    l = x0 + col * (cw + gx); t = y0 + row * (ch + gy)
    rrect(s, l, t, cw, ch)
    circle(s, l + 0.34, t + (ch - 0.26) / 2, 0.26, AMBER)
    line_text(s, l + 0.86, t, cw - 1.1, ch, item, 15, INK, True, BODY, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 14)
notes(s, "Co dalej: więcej danych, więcej miast, inne modele do porównania, automatyzacja i baza wyników. "
         "Zgodne z sekcją „dalszy rozwój” w sprawozdaniu. (~1 min)")

# ===================== SLIDE 15 — CLOSING =====================
s = slide(NAVY)
vdash(s, 2.6, 1.4, 4.7, AMBER, 2.2)
line_text(s, 1.4, 2.55, 11.0, 1.1, "Dziękujemy za uwagę", 46, WHITE, True, HEAD)
line_text(s, 1.45, 3.8, 10.0, 0.7, "Pytania?", 26, AMBER, True, HEAD)
line_text(s, 1.45, 5.6, 11.0, 0.5, "AirSense Weather AI  ·  Grupa 5", 16, ICEBL)
line_text(s, 1.45, 6.1, 11.0, 0.5, "streamlit run app.py", 14, WHITE, False, "Consolas")
notes(s, "Podziękowanie i otwarcie na pytania. (~30 s)")

prs.save(OUT)
print("Zapisano:", OUT, "| slajdów:", len(prs.slides._sldIdLst))
